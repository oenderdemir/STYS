"""STYS Kurumsal Sunum - PPTX Generator"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import re

# ── Türkçe Karakter Düzeltme ──────────────────────────────────
# Kök bazlı harita: her kök kelimenin başına \b (word boundary) eklenir,
# böylece "ici" → "içi" dönüşümü "yönetici" kelimesini bozmaz.
# Sıralama: uzun → kısa (önce uzun eşleşme, sonra kısa).
_TR_ROOTS = [
    # --- ö ---
    ("yonetilebilirlik", "yönetilebilirlik"),
    ("yonetebilir", "yönetebilir"),
    ("yonetiliyor", "yönetiliyor"),
    ("yonetilir", "yönetilir"),
    ("yonetmeyi", "yönetmeyi"),
    ("yonetici", "yönetici"),
    ("yonetim", "yönetim"),
    ("yonetir", "yönetir"),
    ("odemesiz", "ödemesiz"),
    ("odenmis", "ödenmiş"),
    ("olusturmama", "oluşturmama"),
    ("olustur", "oluştur"),
    ("gosterilebilir", "gösterilebilir"),
    ("gosterimi", "gösterimi"),
    ("gorunurlugu", "görünürlüğü"),
    ("gorunurluk", "görünürlük"),
    ("gorunumu", "görünümü"),
    ("gorunum", "görünüm"),
    ("gorunur", "görünür"),
    ("gorebilir", "görebilir"),
    ("goremez", "göremez"),
    ("gorulur", "görülür"),
    ("gorevli", "görevli"),
    ("gorev", "görev"),
    ("gorur", "görür"),
    ("gore", "göre"),
    ("donemsel", "dönemsel"),
    ("donem", "dönem"),
    ("doner", "döner"),
    ("cozulur", "çözülür"),
    ("cozum", "çözüm"),
    ("cozer", "çözer"),
    ("onizleme", "önizleme"),
    ("oncesi", "öncesi"),
    ("oneri", "öneri"),
    ("ozel", "özel"),
    ("ozet", "özet"),
    ("odeme", "ödeme"),
    # --- ü ---
    ("sureclerinden", "süreçlerinden"),
    ("sureclerini", "süreçlerini"),
    ("surecleri", "süreçleri"),
    ("sureci", "süreci"),
    ("surec", "süreç"),
    ("suresi", "süresi"),
    ("sure", "süre"),
    ("ucretlendirme", "ücretlendirme"),
    ("ucretsiz", "ücretsiz"),
    ("ucretli", "ücretli"),
    ("ucretler", "ücretler"),
    ("ucret", "ücret"),
    ("uretimi", "üretimi"),
    ("uretim", "üretim"),
    ("tuketim", "tüketim"),
    ("dusurur", "düşürür"),
    ("dusuyorsa", "düşüyorsa"),
    ("dusen", "düşen"),
    ("duser", "düşer"),
    ("butunlesik", "bütünleşik"),
    ("butunlugu", "bütünlüğü"),
    ("buyumeye", "büyümeye"),
    ("buyutur", "büyütür"),
    ("buyur", "büyür"),
    ("guvenlik", "güvenlik"),
    ("guvence", "güvence"),
    ("guncelleme", "güncelleme"),
    ("guncel", "güncel"),
    ("musteri", "müşteri"),
    ("musait", "müsait"),
    ("menuyu", "menüyü"),
    ("menusunu", "menüsünü"),
    ("menu", "menü"),
    ("moduler", "modüler"),
    ("moduller", "modüller"),
    ("duzeltme", "düzeltme"),
    ("uzerindeki", "üzerindeki"),
    ("uzerinden", "üzerinden"),
    ("uzerinde", "üzerinde"),
    ("surum", "sürüm"),
    ("yuzlerce", "yüzlerce"),
    ("urun", "ürün"),
    # --- ş ---
    ("basvurularini", "başvurularını"),
    ("basvurularim", "başvurularım"),
    ("basvuru", "başvuru"),
    ("baslangic", "başlangıç"),
    ("baslik", "başlık"),
    ("teshisi", "teşhisi"),
    ("teshis", "teşhis"),
    ("islemlerle", "işlemlerle"),
    ("islemleri", "işlemleri"),
    ("isleyisi", "işleyişi"),
    ("isletmeciligi", "işletmeciliği"),
    ("isletmesi", "işletmesi"),
    ("isletme", "işletme"),
    ("islemi", "işlemi"),
    ("islem", "işlem"),
    ("paylasimli", "paylaşımlı"),
    ("paylasim", "paylaşım"),
    ("yerlestirme", "yerleştirme"),
    ("yerlesim", "yerleşim"),
    ("degisiklikler", "değişiklikler"),
    ("degisiklik", "değişiklik"),
    ("degisimi", "değişimi"),
    ("degerlendirmesi", "değerlendirmesi"),
    ("degerlendirip", "değerlendirip"),
    ("degerlendirilir", "değerlendirilir"),
    ("degerlendirme", "değerlendirme"),
    ("genisletilebilir", "genişletilebilir"),
    ("genislemeye", "genişlemeye"),
    ("kisiye", "kişiye"),
    ("kisinin", "kişinin"),
    ("kisi", "kişi"),
    ("iliskilendirme", "ilişkilendirme"),
    ("iliskileri", "ilişkileri"),
    ("iliskisi", "ilişkisi"),
    ("erisebilir", "erişebilir"),
    ("erisim", "erişim"),
    ("disindaki", "dışındaki"),
    ("disi", "dışı"),
    ("kosul", "koşul"),
    ("calismasi", "çalışması"),
    ("calisma", "çalışma"),
    ("calis", "çalış"),
    ("eslesme", "eşleşme"),
    ("esleme", "eşleme"),
    ("gercegine", "gerçeğine"),
    ("gercegi", "gerçeği"),
    ("gercekleri", "gerçekleri"),
    ("gercek", "gerçek"),
    ("siparisleri", "siparişleri"),
    ("siparis", "sipariş"),
    ("secenekleri", "seçenekleri"),
    ("secenek", "seçenek"),
    ("secimi", "seçimi"),
    ("secim", "seçim"),
    ("gecmis", "geçmiş"),
    ("gecis", "geçiş"),
    ("gec ", "geç "),
    ("acilabilen", "açılabilen"),
    ("acilan", "açılan"),
    ("acilir", "açılır"),
    ("acik", "açık"),
    ("icinde", "içinde"),
    ("icin", "için"),
    ("ici", "içi"),
    ("tusluk", "tuşluk"),
    ("sablon", "şablon"),
    # --- ğ ---
    ("daginik", "dağınık"),
    ("sagladigi", "sağladığı"),
    ("saglanan", "sağlanan"),
    ("saglar", "sağlar"),
    ("sagla", "sağla"),
    ("baglantili", "bağlantılı"),
    ("baglan", "bağlan"),
    ("dogrulama", "doğrulama"),
    ("dogrula", "doğrula"),
    ("degistirdi", "değiştirdi"),
    ("deger", "değer"),
    # --- ı ---
    ("karsilastirmali", "karşılaştırmalı"),
    ("karsilastirma", "karşılaştırma"),
    ("katilmayanlarin", "katılmayanların"),
    ("katilimci", "katılımcı"),
    ("katilim", "katılım"),
    ("kisitliysa", "kısıtlıysa"),
    ("kisalir", "kısalır"),
    ("kisalt", "kısalt"),
    ("hazirlaniy", "hazırlanıy"),
    ("hazirlama", "hazırlama"),
    ("hazirlik", "hazırlık"),
    ("hazir", "hazır"),
    ("aninda", "anında"),
    ("anlik", "anlık"),
    ("tanitim", "tanıtım"),
    ("tanimlanabilir", "tanımlanabilir"),
    ("tanimlama", "tanımlama"),
    ("kirilimi", "kırılımı"),
    ("kirilim", "kırılım"),
    ("oranini", "oranını"),
    ("orani", "oranı"),
    ("arttirir", "arttırır"),
    ("artiyorsa", "artıyorsa"),
    ("kaldirir", "kaldırır"),
    ("siralama", "sıralama"),
    ("sinirlamasi", "sınırlaması"),
    ("sinirlamali", "sınırlamalı"),
    ("sinirlama", "sınırlama"),
    ("siniri", "sınırı"),
    ("sinir", "sınır"),
    ("plansiz", "plansız"),
    ("kontrollu", "kontrollü"),
    ("kontrolu", "kontrolü"),
    ("uyarilari", "uyarıları"),
    ("uyarilar", "uyarılar"),
    ("uyari", "uyarı"),
    ("uyumlulugu", "uyumluluğu"),
    ("yukunu", "yükünü"),
    ("yuku", "yükü"),
    ("yuk", "yük"),
    ("parcali", "parçalı"),
    ("varsayilan", "varsayılan"),
    ("hesaplanmasi", "hesaplanması"),
    ("hesabi", "hesabı"),
    ("kapanis", "kapanış"),
    ("kapali", "kapalı"),
    ("cikarildi", "çıkarıldı"),
    ("ciktilari", "çıktıları"),
    ("cikti", "çıktı"),
    ("cikis", "çıkış"),
    ("coklu", "çoklu"),
    ("cocuk", "çocuk"),
    ("hizli", "hızlı"),
    ("hizla", "hızla"),
    ("giris", "giriş"),
    ("kolaylasir", "kolaylaşır"),
    ("farkindalik", "farkındalık"),
    ("farkli", "farklı"),
    ("ayrildi", "ayrıldı"),
    ("sonrasi", "sonrası"),
    ("operasyonlari", "operasyonları"),
    ("zorlasiyor", "zorlaşıyor"),
    ("arizali", "arızalı"),
    ("ariza", "arıza"),
    ("vazgecme", "vazgeçme"),
    ("iptal", "iptal"),
    ("kidem", "kıdem"),
    ("detayli", "detaylı"),
    ("aktarim", "aktarım"),
    ("goc", "göç"),
    ("ihtiyaclarinin", "ihtiyaçlarının"),
    ("ihtiyaclarina", "ihtiyaçlarına"),
    ("tabanli", "tabanlı"),
    ("destegi", "desteği"),
    ("atamasi", "ataması"),
    ("cevaplanir", "cevaplanır"),
    ("izlenebilirlik", "izlenebilirlik"),
    ("yazilim", "yazılım"),
    ("gruplari", "grupları"),
    ("programlari", "programları"),
    ("atamalari", "atamaları"),
    ("kurallari", "kuralları"),
    ("senaryolari", "senaryoları"),
    ("misafirleri", "misafirleri"),
    ("kurumda", "kurumda"),
    ("bazli", "bazlı"),
    ("ekrani", "ekranı"),
    ("odakli", "odaklı"),
    ("ayni", "aynı"),
    ("kilar", "kılar"),
    ("sonuc", "sonuç"),
    ("zayiflar", "zayıflar"),
    ("kullanici", "kullanıcı"),
    ("alani", "alanı"),
    ("gorunumu", "görünümü"),
    # --- Cümle başı büyük harf özel (elle eklenen) ---
    ("Uc uca", "Uç uca"), ("uc uca", "uç uca"),
    ("Is ", "İş "), ("is ", "iş "),
    ("Cok ", "Çok "), ("cok ", "çok "),
]

# Türkçe büyük harf eşleme (Python default upper() Türkçe İ/I sorununu çözmez)
def _tr_upper(ch):
    if ch == 'i': return 'İ'
    if ch == 'ı': return 'I'
    if ch == 'ö': return 'Ö'
    if ch == 'ü': return 'Ü'
    if ch == 'ş': return 'Ş'
    if ch == 'ç': return 'Ç'
    if ch == 'ğ': return 'Ğ'
    return ch.upper()

# Derlenmiş regex listesi: her kök için \b ile word boundary
_TR_COMPILED = []
for _ascii, _tr in _TR_ROOTS:
    # Büyük harf versiyonunu da ekle (ilk harf büyük)
    _cap_a = _tr_upper(_ascii[0]) + _ascii[1:]
    _cap_t = _tr_upper(_tr[0]) + _tr[1:]
    # Önce büyük harfli, sonra küçük harfli
    if _cap_a != _ascii:
        _TR_COMPILED.append((re.compile(r'\b' + re.escape(_cap_a)), _cap_t))
    _TR_COMPILED.append((re.compile(r'\b' + re.escape(_ascii)), _tr))


def fix_turkish(text):
    """Kök bazlı, word-boundary korumalı Türkçe karakter düzeltmesi."""
    result = text
    for pattern, replacement in _TR_COMPILED:
        result = pattern.sub(replacement, result)
    return result

# ── Color Palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # very dark navy
ACCENT_BLUE   = RGBColor(0x00, 0x7A, 0xCC)   # primary blue
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xA0)   # teal accent
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)   # warm orange
LIGHT_BG      = RGBColor(0xF5, 0xF7, 0xFA)   # very light gray bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_GRAY     = RGBColor(0x5A, 0x5A, 0x6E)
TEXT_LIGHT    = RGBColor(0xCC, 0xCC, 0xDD)
CARD_BG       = RGBColor(0x16, 0x21, 0x3A)   # card on dark
STRIPE_BLUE   = RGBColor(0x00, 0x5F, 0xA3)
GRADIENT_MID  = RGBColor(0x0A, 0x2E, 0x5C)
LIGHT_CARD    = RGBColor(0xE8, 0xEE, 0xF6)
SUCCESS_GREEN = RGBColor(0x27, 0xAE, 0x60)
WARN_AMBER    = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


# ── Helper functions ───────────────────────────────────────────

def _add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _tf(shape):
    """Return textframe, cleared."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    return tf


def _add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def _set_para(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, spacing_after=Pt(6), spacing_before=Pt(0), font_name="Segoe UI"):
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "" else tf.add_paragraph()
    p.text = fix_turkish(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = spacing_after
    p.space_before = spacing_before
    return p


def _add_bullet_list(tf, items, size=16, color=WHITE, indent_level=0, bullet_char="\u2022", spacing=Pt(4)):
    for item in items:
        p = tf.add_paragraph()
        p.text = fix_turkish(f"{bullet_char}  {item}")
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.space_before = Pt(0)
        p.level = indent_level


def _slide_number(slide, num, total=20):
    tb = _add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4))
    _set_para(tb.text_frame, f"{num} / {total}", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def _dark_slide_header(slide, title, subtitle=None, slide_num=None):
    """Standard dark slide with title area."""
    _add_bg(slide, DARK_BG)
    # Top accent line
    _add_accent_bar(slide, Inches(0), Inches(0), SW, Inches(0.06), ACCENT_BLUE)
    # Left accent bar
    _add_accent_bar(slide, Inches(0.6), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_TEAL)
    # Title
    tb = _add_text_box(slide, Inches(0.85), Inches(0.6), Inches(11), Inches(0.75))
    _set_para(tb.text_frame, title, size=30, bold=True, color=WHITE)
    if subtitle:
        _set_para(tb.text_frame, subtitle, size=16, color=TEXT_LIGHT, spacing_before=Pt(2))
    if slide_num:
        _slide_number(slide, slide_num)


def _content_card(slide, left, top, width, height, color=CARD_BG):
    return _add_shape(slide, left, top, width, height, color, radius=True)


# ── SLIDE 1: COVER ────────────────────────────────────────────
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, DARK_BG)

    # Decorative shapes
    _add_accent_bar(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
    _add_accent_bar(slide, Inches(0), Inches(7.42), SW, Inches(0.08), ACCENT_TEAL)

    # Side accent
    _add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SH, ACCENT_BLUE)

    # Large decorative circle (top right)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x0A, 0x2E, 0x5C)
    circle.line.fill.background()
    circle.shadow.inherit = False
    circle.rotation = 0

    # Small decorative circle
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(4.5), Inches(3.5), Inches(3.5))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x12, 0x25, 0x45)
    c2.line.fill.background()
    c2.shadow.inherit = False

    # Main title
    tb = _add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2))
    _set_para(tb.text_frame, "STYS", size=64, bold=True, color=ACCENT_TEAL, font_name="Segoe UI Black")

    # Subtitle line
    tb2 = _add_text_box(slide, Inches(1.2), Inches(3.1), Inches(8.5), Inches(1.0))
    _set_para(tb2.text_frame, "Kurumsal Konaklama, Tesis ve Restoran", size=28, bold=True, color=WHITE)
    _set_para(tb2.text_frame, "Operasyon Yonetim Platformu", size=28, bold=True, color=WHITE)

    # Description
    tb3 = _add_text_box(slide, Inches(1.2), Inches(4.5), Inches(8), Inches(1.0))
    _set_para(tb3.text_frame, "Rezervasyon, kamp, restoran, odeme, temizlik ve yetki yonetimini", size=16, color=TEXT_LIGHT)
    _set_para(tb3.text_frame, "tek merkezde toplayan butunlesik dijital cozum", size=16, color=TEXT_LIGHT)

    # Accent line under title
    _add_accent_bar(slide, Inches(1.2), Inches(4.3), Inches(4), Inches(0.05), ACCENT_TEAL)

    # Date
    tb4 = _add_text_box(slide, Inches(1.2), Inches(6.2), Inches(5), Inches(0.4))
    _set_para(tb4.text_frame, "Nisan 2026  |  Kurumsal Tanitim Sunumu", size=13, color=TEXT_GRAY)


# ── SLIDE 2: PROBLEM ──────────────────────────────────────────
def slide_problem():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Konaklama ve tesis operasyonlari neden zorlasiyor?", slide_num=2)

    # Main message
    tb = _add_text_box(slide, Inches(0.9), Inches(1.6), Inches(7), Inches(0.5))
    _set_para(tb.text_frame, "Bir cok kurumda su alanlar birbirinden kopuk yonetiliyor:", size=17, color=TEXT_LIGHT)

    problems = [
        ("Rezervasyon alma ve senaryo olusturma", ""),
        ("Oda ve yatak yerlestirme", ""),
        ("Check-in / check-out operasyonu", ""),
        ("Odeme, ek hizmet ve indirim yonetimi", ""),
        ("Temizlik ve oda hazirlik surecleri", ""),
        ("Arizali oda ve operasyonel degisiklikler", ""),
        ("Kamp / toplu konaklama basvuru ve tahsis", ""),
        ("Restoran, siparis ve masa yonetimi", ""),
        ("Bildirim ve anlik operasyonel takip", ""),
    ]

    # Left column - problem list in a card
    card = _content_card(slide, Inches(0.7), Inches(2.2), Inches(7.2), Inches(4.5))
    ctf = _tf(card)
    ctf.paragraphs[0].space_before = Pt(10)
    for prob, _ in problems:
        p = ctf.add_paragraph()
        p.text = fix_turkish(f"    \u2716   {prob}")
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
        p.font.name = "Segoe UI"
        p.space_after = Pt(5)

    # Right side - warning card
    warn_card = _content_card(slide, Inches(8.5), Inches(2.2), Inches(4.2), Inches(4.5), RGBColor(0x2D, 0x15, 0x15))
    _add_accent_bar(slide, Inches(8.5), Inches(2.2), Inches(0.08), Inches(4.5), ACCENT_ORANGE)
    wtf = _tf(warn_card)
    wtf.paragraphs[0].space_before = Pt(20)
    _set_para(wtf, "   Sonuc", size=22, bold=True, color=ACCENT_ORANGE, spacing_after=Pt(20))
    warnings = [
        "Hata orani artar",
        "Yonetim gorunurlugu duser",
        "Personel is yuku gereksiz yere buyur",
        "Denetlenebilirlik zayiflar",
        "Kapasite israf edilir",
    ]
    for w in warnings:
        p = wtf.add_paragraph()
        p.text = fix_turkish(f"    \u26A0  {w}")
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xCC, 0x80)
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)

    _slide_number(slide, 2)


# ── SLIDE 3: SOLUTION ─────────────────────────────────────────
def slide_solution():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "STYS ne cozer?", slide_num=3)

    tb = _add_text_box(slide, Inches(0.9), Inches(1.6), Inches(11), Inches(0.6))
    _set_para(tb.text_frame,
              "Konaklama operasyonundan restoran yonetimine, kamp sureclerinden finansal takibe kadar tum isleyisi uc uca yonetir.",
              size=16, color=TEXT_LIGHT)

    modules = [
        ("\U0001F4C5", "Akilli Rezervasyon\nve Senaryo Motoru", ACCENT_BLUE),
        ("\U0001F6CF", "Kisi, Oda ve Yatak\nBazli Planlama", RGBColor(0x00, 0x96, 0xC7)),
        ("\U0001F6AA", "Check-in /\nCheck-out", RGBColor(0x00, 0xB4, 0xD8)),
        ("\U0001F4B3", "Odeme, Ek Hizmet\nve Indirim", ACCENT_TEAL),
        ("\U0001F9F9", "Oda Temizlik ve\nHazirlik Surecleri", RGBColor(0x2D, 0xC6, 0x53)),
        ("\U0001F527", "Arizali Oda ve\nOda Degisimi", RGBColor(0x48, 0xBF, 0x91)),
        ("\U0001F3D5", "Kamp Yonetimi\n(Basvuru-Tahsis-Iade)", RGBColor(0x84, 0x5E, 0xC2)),
        ("\U0001F37D", "Restoran ve\nSiparis Yonetimi", RGBColor(0xE0, 0x77, 0x00)),
        ("\U0001F4F1", "Musteri\nDijital Menu", WARN_AMBER),
        ("\U0001F514", "Bildirim ve\nAnlik Uyarilar", RGBColor(0xE7, 0x4C, 0x3C)),
        ("\U0001F4CA", "Dashboard, KPI\nve Raporlama", RGBColor(0x34, 0x98, 0xDB)),
        ("\U0001F512", "Rol, Yetki ve\nErisim Teshisi", RGBColor(0x95, 0xA5, 0xA6)),
    ]

    cols = 4
    card_w = Inches(2.7)
    card_h = Inches(1.35)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = Inches(0.7)
    start_y = Inches(2.5)

    for i, (icon, label, color) in enumerate(modules):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        # top color stripe
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), color)

        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.2))
        _set_para(tb.text_frame, icon, size=24, color=color, spacing_after=Pt(4))
        _set_para(tb.text_frame, label, size=13, bold=True, color=WHITE, spacing_after=Pt(0))


# ── SLIDE 4: TARGET ───────────────────────────────────────────
def slide_target():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Hedef Kurum Tipleri", "Kimler icin uygun?", slide_num=4)

    targets = [
        ("\U0001F3DB", "Kamu Misafirhaneleri", "Farkli sehirlerden gelen personel ve misafirler icin planli konaklama"),
        ("\U0001F393", "Universite Konukevleri", "Akademisyen, arastirmaci ve davetli konuklar icin yonetim"),
        ("\U0001F3D6", "Sosyal Tesisler", "Dinlenme merkezleri ve sosyal tesis isletmeleri"),
        ("\U0001F3E2", "Kurum Lojmanlari", "Gecici konaklama ve lojman yonetimi"),
        ("\U0001F3E8", "Butik Oteller", "Pansiyon ve ozel konaklama tesisleri"),
        ("\U0001F3D5", "Kamp Tesisleri", "Toplu konaklama ve kamp organizasyonlari"),
        ("\U0001F3EB", "Kampus Merkezleri", "Misafir kabul ve konuk koordinasyonu"),
        ("\U0001F37D", "Restoran Tesisleri", "Restoran isletmesi bulunan tum tesis tipleri"),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.0)

    for i, (icon, title, desc) in enumerate(targets):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        tb = _add_text_box(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.25))
        _set_para(tb.text_frame, icon, size=30, color=ACCENT_TEAL, spacing_after=Pt(6))
        _set_para(tb.text_frame, title, size=15, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=11, color=TEXT_LIGHT, spacing_after=Pt(0))


# ── SLIDE 5: VALUE PROPOSITION ─────────────────────────────────
def slide_value():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Kurumlara Saglanan Katma Deger", slide_num=5)

    values = [
        ("\U0001F517", "Tek Sistem", "Rezervasyondan operasyona,\nrestorandan kampa entegre cozum"),
        ("\U0001F4C8", "Verimli Kapasite", "Izlenebilir ve optimize\nedilebilir kapasite kullanimi"),
        ("\u2699\uFE0F", "Minimum Is Yuku", "Manuel takip ve kirtasiye\nislemlerinin ortadan kalkmasi"),
        ("\U0001F50D", "Denetlenebilirlik", "Raporlanabilir ve\nizlenebilir surecler"),
        ("\U0001F4B0", "Finansal Gorunurluk", "Gelir takibi ve\nodeme kirilim raporlari"),
        ("\U0001F512", "Yetki Kontrolu", "Merkezi kapsam ve\nerisim yonetimi"),
        ("\U0001F514", "Anlik Bildirim", "Operasyonel farkindalik\nve uyari sistemi"),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(2.1)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(1.9)

    for i, (icon, title, desc) in enumerate(values):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), ACCENT_TEAL)
        tb = _add_text_box(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.3))
        _set_para(tb.text_frame, icon, size=32, color=ACCENT_TEAL, spacing_after=Pt(6))
        _set_para(tb.text_frame, title, size=16, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)


# ── SLIDE 6: SMART RESERVATION ────────────────────────────────
def slide_reservation():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Rezervasyon ve Senaryo Motoru", "Akilli Rezervasyon", slide_num=6)

    # Left - Features card
    card = _content_card(slide, Inches(0.7), Inches(1.9), Inches(7.5), Inches(5.0))
    _add_accent_bar(slide, Inches(0.7), Inches(1.9), Inches(0.06), Inches(5.0), ACCENT_BLUE)
    ctf = _tf(card)
    ctf.paragraphs[0].space_before = Pt(12)

    features = [
        "Musaitlik kontrolu ve oda kapasite analizi",
        "Minimum gece ve stop-sale kurallari",
        "Segmentli veya tek parca alternatif senaryolar",
        "Gereksiz segment olusturmama optimizasyonu",
        "Kisi sayisi ve kapasite uyumlulugu",
        "Paylasimli oda mantigi ve cinsiyet uyumu",
        "Konaklama tipi ve misafir tipi bazli planlama",
        "Konaklama hakki / kontenjan tuketim takibi",
        "Indirim kurali yonetimi (otomatik ve manuel)",
    ]

    for f in features:
        p = ctf.add_paragraph()
        p.text = fix_turkish(f"     \u2713   {f}")
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    # Right - benefit highlight
    benefit = _content_card(slide, Inches(8.7), Inches(1.9), Inches(4.0), Inches(5.0), RGBColor(0x0A, 0x3D, 0x2A))
    _add_accent_bar(slide, Inches(8.7), Inches(1.9), Inches(4.0), Inches(0.06), ACCENT_TEAL)
    btf = _tf(benefit)
    btf.paragraphs[0].space_before = Pt(20)
    _set_para(btf, "   \U0001F4A1", size=36, spacing_after=Pt(10))
    _set_para(btf, "   Kurum Faydasi", size=18, bold=True, color=ACCENT_TEAL, spacing_after=Pt(16))
    _set_para(btf, "   Resepsiyon veya operasyon\n   personeli elde hesap yapmak\n   zorunda kalmaz.", size=14, color=WHITE, spacing_after=Pt(12))
    _set_para(btf, "   Sistem en uygun konaklama\n   senaryolarini otomatik\n   olusturur.", size=14, color=TEXT_LIGHT)


# ── SLIDE 7: GUEST PLAN ───────────────────────────────────────
def slide_guest_plan():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Kisi Bazli ve Yatak Bazli Yerlesim", "Konaklayan Plani", slide_num=7)

    features = [
        ("\U0001F464", "Bireysel Plan", "Her konaklayan icin ayri\nkonaklama plani olusturma"),
        ("\U0001F6CF", "Yatak Atamasi", "Paylasimli odada yatak\nbazli atama ve takip"),
        ("\u2642\u2640", "Cinsiyet Uyumu", "Paylasimli odalarda\ncinsiyet uyumlulugu kontrolu"),
        ("\U0001F4CB", "Durum Takibi", "Gelen / Bekleyen / Gelmedi\n/ Ayrildi durumlari"),
        ("\U0001F4B3", "Hizmet Baglama", "Konaklayan bazinda ek hizmet\nve odeme iliskilendirme"),
    ]

    card_w = Inches(2.25)
    card_h = Inches(3.2)
    gap = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.2)

    for i, (icon, title, desc) in enumerate(features):
        x = start_x + i * (card_w + gap)
        card = _content_card(slide, x, start_y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, start_y, card_w, Inches(0.05), ACCENT_BLUE)
        tb = _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.3))
        _set_para(tb.text_frame, icon, size=36, color=ACCENT_BLUE, spacing_after=Pt(8), align=PP_ALIGN.CENTER)
        _set_para(tb.text_frame, title, size=15, bold=True, color=WHITE, spacing_after=Pt(10), align=PP_ALIGN.CENTER)
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Bottom benefit
    tb = _add_text_box(slide, Inches(0.9), Inches(5.8), Inches(11), Inches(0.6))
    _set_para(tb.text_frame, "\U0001F4A1  Gercek hayattaki konaklama plani sistemde net olarak gorulur; kim nerede, ne zaman kalacak sorusu aninda cevaplanir.",
              size=14, color=ACCENT_TEAL)


# ── SLIDE 8: CHECK-IN / CHECK-OUT ─────────────────────────────
def slide_checkin():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Kontrollu Giris ve Cikis Operasyonu", "Check-in / Check-out", slide_num=8)

    items = [
        ("\U0001F6AB", "Plansiz check-in engeli", "Plan olmadan check-in yapilamaz"),
        ("\u26A0\uFE0F", "Oda hazirlik kontrolu", "Hazir olmayan oda icin uyari veya engel"),
        ("\U0001F4B3", "Odeme kontrolu", "Eksik odeme varsa check-out engeli"),
        ("\U0001F464", "Fiili operasyon", "Gelen misafirler uzerinden yonetim"),
        ("\u23F0", "Tesis bazli saatler", "Her tesis icin farkli giris/cikis saati"),
        ("\U0001F4C5", "Otomatik gun hesabi", "Tesis kuralina gore otomatik hesaplama"),
    ]

    cols = 3
    card_w = Inches(3.8)
    card_h = Inches(1.6)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.1)

    for i, (icon, title, desc) in enumerate(items):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.25))
        _set_para(tb.text_frame, f"{icon}  {title}", size=15, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)


# ── SLIDE 9: PAYMENT ──────────────────────────────────────────
def slide_payment():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Gelir, Odeme ve Hizmet Takibi", "Odeme ve Ek Hizmetler", slide_num=9)

    # Left - payment features
    card = _content_card(slide, Inches(0.7), Inches(1.9), Inches(6.5), Inches(4.6))
    _add_accent_bar(slide, Inches(0.7), Inches(1.9), Inches(0.06), Inches(4.6), ACCENT_TEAL)
    ctf = _tf(card)
    ctf.paragraphs[0].space_before = Pt(12)

    feats = [
        "Parcali odeme (kademeli tahsilat)",
        "Nakit veya kart ile odeme",
        "Ek hizmet satisi ve kisi bazli hizmet yazma",
        "Varsayilan fiyat + serbest fiyat override",
        "Indirim kurali yonetimi",
        "Konaklama hakki tuketim kirilimi",
    ]
    for f in feats:
        p = ctf.add_paragraph()
        p.text = fix_turkish(f"     \u2713   {f}")
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(7)

    # Right - example services
    svc_card = _content_card(slide, Inches(7.7), Inches(1.9), Inches(5.0), Inches(4.6), RGBColor(0x1A, 0x2A, 0x1A))
    _add_accent_bar(slide, Inches(7.7), Inches(1.9), Inches(5.0), Inches(0.05), ACCENT_ORANGE)
    stf = _tf(svc_card)
    stf.paragraphs[0].space_before = Pt(12)
    _set_para(stf, "   Ornek Ek Hizmetler", size=17, bold=True, color=ACCENT_ORANGE, spacing_after=Pt(14))

    services = [
        "\U0001F45F  Ayakkabi boyama",
        "\U0001F455  Kurutemizleme",
        "\u2615  Odaya kahvalti",
        "\U0001F698  Transfer",
        "\u23F0  Gec check-out",
        "\u2728  Ozel istekler",
    ]
    for s in services:
        p = stf.add_paragraph()
        p.text = fix_turkish(f"     {s}")
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)


# ── SLIDE 10: HOUSEKEEPING ────────────────────────────────────
def slide_housekeeping():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Housekeeping ve Oda Hazirlik Sureci", "Oda Temizlik ve Hazirlik", slide_num=10)

    steps = [
        ("\U0001F6AA", "Check-out", "Otomatik kirli\ndurum tanimlama"),
        ("\u27A1\uFE0F", "", ""),
        ("\U0001F9F9", "Temizlik", "Gorevliye is dusmesi\nTemizleniyor durumu"),
        ("\u27A1\uFE0F", "", ""),
        ("\u2705", "Hazir", "Hazir durumuna\ngecis"),
        ("\u27A1\uFE0F", "", ""),
        ("\U0001F6AA", "Check-in", "Hazirlik kontrolu\nsonrasi giris"),
    ]

    x = Inches(0.5)
    y = Inches(2.8)

    for i, (icon, title, desc) in enumerate(steps):
        if icon == "\u27A1\uFE0F":
            # Arrow
            tb = _add_text_box(slide, x, y + Inches(0.4), Inches(0.8), Inches(0.6))
            _set_para(tb.text_frame, "\u27A1", size=28, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
            x += Inches(0.8)
        else:
            card_w = Inches(2.4)
            card_h = Inches(2.5)
            card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
            _add_accent_bar(slide, x, y, card_w, Inches(0.05), ACCENT_BLUE)
            tb = _add_text_box(slide, x + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), card_h - Inches(0.2))
            _set_para(tb.text_frame, icon, size=36, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, spacing_after=Pt(8))
            _set_para(tb.text_frame, title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, spacing_after=Pt(8))
            _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
            x += card_w + Inches(0.1)


# ── SLIDE 11: ROOM MAINTENANCE ────────────────────────────────
def slide_maintenance():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Operasyonel Kesintilere Hazir Yapi", "Arizali Oda ve Oda Degisimi", slide_num=11)

    items = [
        ("\U0001F527", "Ariza Kaydi", "Arizali veya kullanima kapali\noda kaydi olusturma"),
        ("\U0001F50D", "Etki Tespiti", "Etkilenen rezervasyonlari\notomatik tespit etme"),
        ("\U0001F4A1", "Alternatif Oneri", "Uygun alternatif oda\nonerisi ve atama"),
        ("\U0001F504", "Oda Degisimi", "Check-in oncesi ve sonrasi\noda degisimi destegi"),
        ("\U0001F4B0", "Fiyat Koruma", "Fiyat artiyorsa koruma,\ndusuyorsa indirim uygulama"),
    ]

    card_w = Inches(2.25)
    card_h = Inches(3.0)
    gap = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.2)

    for i, (icon, title, desc) in enumerate(items):
        x = start_x + i * (card_w + gap)
        card = _content_card(slide, x, start_y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, start_y, card_w, Inches(0.05), ACCENT_ORANGE)
        tb = _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.3))
        _set_para(tb.text_frame, icon, size=36, color=ACCENT_ORANGE, spacing_after=Pt(8), align=PP_ALIGN.CENTER)
        _set_para(tb.text_frame, title, size=15, bold=True, color=WHITE, spacing_after=Pt(10), align=PP_ALIGN.CENTER)
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ── SLIDE 12: CAMP MANAGEMENT ─────────────────────────────────
def slide_camp():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Toplu Konaklama ve Kamp Surec Yonetimi", "Kamp Yonetimi", slide_num=12)

    items = [
        ("\U0001F4C5", "Program ve Donem", "Kamp programlari, donemler\nve tesis atamalari"),
        ("\U0001F4DD", "Online Basvuru", "Katilimci bilgileri, konaklama\nbirimi secimi, ucret onizleme"),
        ("\U0001F3AF", "Puan Bazli Tahsis", "Kidem, gecmis katilim gibi\nkriterlere gore siralama"),
        ("\U0001F504", "Otomatik Rezervasyon", "Tahsisten otomatik konaklama\nrezervasyonu + oda atamasi"),
        ("\U0001F4B0", "Ucret ve Tarife", "Kamu/diger tarife, cocuk\nyas kurallari, birim ucret"),
        ("\U0001F4B8", "Iade ve Iptal", "Vazgecme suresi, kesinti\nhesabi, mazeret degerlendirme"),
        ("\u274C", "No-Show Yonetimi", "Katilmayanlarin\notomatik iptali"),
        ("\U0001F464", "Basvurularim", "Self-servis basvuru\ntakip ekrani"),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(1.8)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.0)

    for i, (icon, title, desc) in enumerate(items):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), RGBColor(0x84, 0x5E, 0xC2))
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3), card_h - Inches(0.2))
        _set_para(tb.text_frame, f"{icon}  {title}", size=14, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)

    # Bottom benefit
    tb = _add_text_box(slide, Inches(0.9), Inches(6.2), Inches(11), Inches(0.5))
    _set_para(tb.text_frame,
              "\U0001F4A1  Yuzlerce basvuruyu adil puanlama ile degerlendirip, tahsis, rezervasyon ve ucretlendirme sureclerini tek tusluk islemlerle tamamlar.",
              size=14, color=RGBColor(0x84, 0x5E, 0xC2))


# ── SLIDE 13: RESTAURANT ──────────────────────────────────────
def slide_restaurant():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Tesis Ici Restoran, Siparis ve Odeme Yonetimi", "Restoran Yonetimi", slide_num=13)

    items = [
        ("\U0001F37D", "Restoran Tanimlama", "Coklu restoran, tesis ve\nisletme alani iliskisi"),
        ("\U0001F4BA", "Masa Yonetimi", "Durum takibi:\nMusait / Rezerve / Kapali"),
        ("\U0001F4DC", "Menu Yonetimi", "Kategori bazli urun, fiyat\nve hazirlama suresi"),
        ("\U0001F5C2", "Kategori Havuzu", "Merkezi kategori sablonu\nve restoran atamasi"),
        ("\U0001F4E6", "Siparis Yonetimi", "Kalem bazli durum:\nBeklemede > Hazir > Servis"),
        ("\U0001F44B", "Garson Servisi", "Dokunmatik uyumlu\nhizli operasyon ekrani"),
        ("\U0001F4B3", "Restoran Odeme", "Nakit/kart, siparis bazli\nodeme ozeti"),
        ("\U0001F4F1", "Dijital Menu", "QR ile acilan, mobil uyumlu\nmusteri menu ekrani"),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(1.8)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.0)

    for i, (icon, title, desc) in enumerate(items):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), ACCENT_ORANGE)
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3), card_h - Inches(0.2))
        _set_para(tb.text_frame, f"{icon}  {title}", size=14, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)

    # Benefit
    tb = _add_text_box(slide, Inches(0.9), Inches(6.2), Inches(11), Inches(0.5))
    _set_para(tb.text_frame,
              "\U0001F4A1  Garsondan yoneticiye, mutfaktan kasaya tek sistem. Konaklama ile ayni platformda butunlesik restoran yonetimi.",
              size=14, color=ACCENT_ORANGE)


# ── SLIDE 14: NOTIFICATION ────────────────────────────────────
def slide_notification():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Anlik Bildirim ve Operasyonel Uyari", "Bildirim Sistemi", slide_num=14)

    items = [
        ("\u26A1", "SignalR Altyapisi", "Gercek zamanli, sunucu tabanli\nanlik bildirim gonderimi"),
        ("\U0001F514", "Operasyonel Uyari", "Kritik olaylarda kullaniciya\naninda haber verme"),
        ("\U0001F7E2\U0001F7E1\U0001F534", "Seviye Bazli", "Bilgi, uyari ve kritik\nseviye ayrimi"),
        ("\U0001F464", "Kisi Bazli", "Kullanici hedefli\nbildirim gonderimi"),
    ]

    card_w = Inches(2.85)
    card_h = Inches(3.0)
    gap = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.2)

    for i, (icon, title, desc) in enumerate(items):
        x = start_x + i * (card_w + gap)
        card = _content_card(slide, x, start_y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, start_y, card_w, Inches(0.05), RGBColor(0xE7, 0x4C, 0x3C))
        tb = _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.3), card_w - Inches(0.3), card_h - Inches(0.4))
        _set_para(tb.text_frame, icon, size=36, color=RGBColor(0xE7, 0x4C, 0x3C), align=PP_ALIGN.CENTER, spacing_after=Pt(10))
        _set_para(tb.text_frame, title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, spacing_after=Pt(10))
        _set_para(tb.text_frame, desc, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ── SLIDE 15: ACCESS DIAGNOSTICS ──────────────────────────────
def slide_access_diag():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Admin ve Destek Ekipleri Icin Hizli Teshis", "Yetki ve Erisim Teshisi", slide_num=15)

    questions = [
        "Kullanici neden menuyu goremez?",
        "Neden listeleyemez veya yeni kayit ekleyemez?",
        "Neden sadece belirli tesiste engellenir?",
        "Modul bazli yetki analiz ekrani",
    ]

    # Questions card
    card = _content_card(slide, Inches(0.7), Inches(2.0), Inches(7.5), Inches(3.8))
    _add_accent_bar(slide, Inches(0.7), Inches(2.0), Inches(0.06), Inches(3.8), ACCENT_BLUE)
    ctf = _tf(card)
    ctf.paragraphs[0].space_before = Pt(20)

    for q in questions:
        p = ctf.add_paragraph()
        p.text = fix_turkish(f"     \u2753   {q}")
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(14)

    # Benefit card
    ben = _content_card(slide, Inches(8.7), Inches(2.0), Inches(4.0), Inches(3.8), RGBColor(0x0A, 0x3D, 0x2A))
    _add_accent_bar(slide, Inches(8.7), Inches(2.0), Inches(4.0), Inches(0.06), SUCCESS_GREEN)
    btf = _tf(ben)
    btf.paragraphs[0].space_before = Pt(20)
    _set_para(btf, "   \u2705", size=36, spacing_after=Pt(12))
    _set_para(btf, "   Kurum Faydasi", size=18, bold=True, color=SUCCESS_GREEN, spacing_after=Pt(14))
    _set_para(btf, "   Destek suresi kisalir.", size=14, color=WHITE, spacing_after=Pt(8))
    _set_para(btf, "   Rol/scope problemleri\n   hizla cozulur.", size=14, color=WHITE, spacing_after=Pt(8))
    _set_para(btf, "   IT ekibinin uzerindeki\n   destek yuku azalir.", size=14, color=TEXT_LIGHT)


# ── SLIDE 16: AUTH ARCHITECTURE ────────────────────────────────
def slide_auth():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Cok Katmanli Yetkilendirme Yapisi", "Yetki Mimarisi", slide_num=16)

    layers = [
        ("\U0001F465", "Rol Bazli", "Admin, Tesis Yonetici,\nBina Yonetici, Resepsiyonist,\nRestoran Yonetici, Garson", ACCENT_BLUE),
        ("\U0001F3E2", "Kapsam Bazli", "Tesis, bina, restoran\nseviyesinde erisim\nsinirlamasi", RGBColor(0x00, 0x96, 0xC7)),
        ("\U0001F3F7", "Marker Roller", "Atanabilir / Atayabilir\nyetki ayrimi\n(ornegin: yonetici atayabilir mi?)", ACCENT_TEAL),
        ("\U0001F4CB", "Menu Gorunurlugu", "Yetkiye gore\ndinamik menu agaci\nolusturma", RGBColor(0x48, 0xBF, 0x91)),
        ("\U0001F46B", "Kullanici Gruplari", "Grup bazli toplu\nyetki atamasi\nve yonetimi", RGBColor(0x84, 0x5E, 0xC2)),
    ]

    card_w = Inches(2.25)
    card_h = Inches(3.5)
    gap = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.0)

    for i, (icon, title, desc, color) in enumerate(layers):
        x = start_x + i * (card_w + gap)
        card = _content_card(slide, x, start_y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, start_y, card_w, Inches(0.06), color)
        tb = _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.3))
        _set_para(tb.text_frame, icon, size=36, color=color, align=PP_ALIGN.CENTER, spacing_after=Pt(8))
        _set_para(tb.text_frame, title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, spacing_after=Pt(10))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Benefit
    tb = _add_text_box(slide, Inches(0.9), Inches(6.0), Inches(11), Inches(0.5))
    _set_para(tb.text_frame,
              "\U0001F4A1  Her kullanici sadece kendi isini gorebilir ve yapabilir. Farkli gorev seviyelerine gore hassas erisim kontrolu.",
              size=14, color=ACCENT_TEAL)


# ── SLIDE 17: DASHBOARD ───────────────────────────────────────
def slide_dashboard():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Anlik Operasyonel Gorunum ve KPI Takibi", "Dashboard ve Raporlama", slide_num=17)

    items = [
        ("\U0001F6AA", "Check-in/out Takibi", "Bugun giris ve cikis\nyapacak misafirler"),
        ("\U0001F3E8", "Kapasite Gosterimi", "Dolu / bos kapasite\nanlik gorunum"),
        ("\U0001F4C8", "KPI Trendleri", "Donemsel karsilastirma\nve trend grafikleri"),
        ("\U0001F4B0", "Gelir Kirilimi", "Baz ucret, indirim ve\nodeme detay raporlari"),
        ("\U0001F3E2", "Coklu Tesis", "Tesis bazli\nkarsilastirmali raporlama"),
        ("\U0001F4C4", "Export", "Excel ve PDF\ncikti destegi"),
    ]

    cols = 3
    card_w = Inches(3.8)
    card_h = Inches(1.7)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.0)

    for i, (icon, title, desc) in enumerate(items):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), RGBColor(0x34, 0x98, 0xDB))
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.25))
        _set_para(tb.text_frame, f"{icon}  {title}", size=15, bold=True, color=WHITE, spacing_after=Pt(6))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)


# ── SLIDE 18: TECHNICAL ───────────────────────────────────────
def slide_technical():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Kurumsal Teknik Guvence", "Teknik Altyapi", slide_num=18)

    stack = [
        ("\u2699\uFE0F", "Backend", "ASP.NET Core (.NET 10)\nEntity Framework Core\nSQL Server", ACCENT_BLUE),
        ("\U0001F5A5", "Frontend", "Angular\n(Standalone Component)\nPrimeNG", RGBColor(0xDD, 0x00, 0x31)),
        ("\u26A1", "Gercek Zamanli", "SignalR ile\nanlik bildirim\naltyapisi", RGBColor(0x84, 0x5E, 0xC2)),
        ("\U0001F3D7", "Mimari", "Repository pattern\nServis katmani\nAudit + Soft-delete", ACCENT_TEAL),
        ("\U0001F512", "Guvenlik", "JWT kimlik dogrulama\nRol/kapsam bazli\nyetkilendirme", SUCCESS_GREEN),
        ("\U0001F4BE", "Veritabani", "Migration bazli\nsurum yonetimi\nSeed data", RGBColor(0x34, 0x98, 0xDB)),
        ("\U0001F680", "Deploy", "Konteyner ve\nkurumsal ortam\nuyumlu", ACCENT_ORANGE),
        ("\U0001F9EA", "Test", "Backend test\naltyapisi\nmevcut", WARN_AMBER),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(2.0)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.0)

    for i, (icon, title, desc, color) in enumerate(stack):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), color)
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), card_h - Inches(0.2))
        _set_para(tb.text_frame, f"{icon}  {title}", size=15, bold=True, color=color, spacing_after=Pt(8))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)


# ── SLIDE 19: WHY STYS ────────────────────────────────────────
def slide_why():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide_header(slide, "Neden Bu Cozum?", "Neden STYS?", slide_num=19)

    reasons = [
        ("\u2705", "Tam Operasyon", "Rezervasyon disindaki tum operasyonu\nda yonetir: restoran, kamp, temizlik, ariza"),
        ("\u2705", "Saha Gercegi", "Paylasimli oda, parcali odeme, no-show\ngibi gercek senaryolari destekler"),
        ("\u2705", "Kurumsal Uyum", "Rol ve scope mantigina\ntam uyumlu yetki yapisi"),
        ("\u2705", "Genisletilebilir", "Denetlenebilir, raporlanabilir\nve moduler yapi"),
        ("\u2705", "Tek Platform", "Konaklama + restoran + kamp\nbutunlugu tek merkezde"),
        ("\u2705", "Mobil Uyumlu", "Dokunmatik operasyon destegi\nve mobil uyumlu ekranlar"),
        ("\u2705", "Hizli Kurulum", "Hazir seed verisi ile\naninda demo ortami"),
    ]

    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(1.8)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.0)

    for i, (icon, title, desc) in enumerate(reasons):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = _content_card(slide, x, y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, y, card_w, Inches(0.05), SUCCESS_GREEN)
        tb = _add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3), card_h - Inches(0.2))
        _set_para(tb.text_frame, f"{icon}  {title}", size=14, bold=True, color=SUCCESS_GREEN, spacing_after=Pt(8))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT)


# ── SLIDE 20: CLOSING ─────────────────────────────────────────
def slide_closing():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BG)

    # Decorative
    _add_accent_bar(slide, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
    _add_accent_bar(slide, Inches(0), Inches(7.42), SW, Inches(0.08), ACCENT_TEAL)
    _add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SH, ACCENT_TEAL)

    # Title
    tb = _add_text_box(slide, Inches(1.2), Inches(1.2), Inches(10), Inches(0.8))
    _set_para(tb.text_frame, "Bir Sonraki Adim", size=38, bold=True, color=WHITE)

    _add_accent_bar(slide, Inches(1.2), Inches(2.1), Inches(3), Inches(0.05), ACCENT_TEAL)

    options = [
        ("\U0001F3AC", "Canli Demo", "Hazir seed verisi ile\naninda gosterilebilir"),
        ("\U0001F3E2", "Pilot Kurulum", "Secili tesis uzerinde\npilot calisma"),
        ("\U0001F504", "Uyarlama", "Kurum ihtiyaclarina gore\nozel gelistirme"),
        ("\U0001F4CA", "Analiz", "Raporlama ve entegrasyon\nihtiyaclarinin tespiti"),
        ("\U0001F4E6", "Veri Goc", "Mevcut sistemlerle\nveri aktarim plani"),
    ]

    card_w = Inches(2.25)
    card_h = Inches(2.8)
    gap = Inches(0.2)
    start_x = Inches(0.65)
    start_y = Inches(2.8)

    for i, (icon, title, desc) in enumerate(options):
        x = start_x + i * (card_w + gap)
        card = _content_card(slide, x, start_y, card_w, card_h, CARD_BG)
        _add_accent_bar(slide, x, start_y, card_w, Inches(0.05), ACCENT_TEAL)
        tb = _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.2), card_w - Inches(0.3), card_h - Inches(0.3))
        _set_para(tb.text_frame, icon, size=36, color=ACCENT_TEAL, align=PP_ALIGN.CENTER, spacing_after=Pt(8))
        _set_para(tb.text_frame, title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, spacing_after=Pt(8))
        _set_para(tb.text_frame, desc, size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Contact placeholder
    tb = _add_text_box(slide, Inches(1.2), Inches(6.3), Inches(10), Inches(0.5))
    _set_para(tb.text_frame, "STYS  |  Kurumsal Konaklama, Tesis ve Restoran Operasyon Yonetim Platformu",
              size=12, color=TEXT_GRAY, align=PP_ALIGN.LEFT)

    _slide_number(slide, 20)


# ── Generate all slides ───────────────────────────────────────
slide_cover()
slide_problem()
slide_solution()
slide_target()
slide_value()
slide_reservation()
slide_guest_plan()
slide_checkin()
slide_payment()
slide_housekeeping()
slide_maintenance()
slide_camp()
slide_restaurant()
slide_notification()
slide_access_diag()
slide_auth()
slide_dashboard()
slide_technical()
slide_why()
slide_closing()

output_path = os.path.join(os.path.dirname(__file__), "STYS_Kurumsal_Sunum.pptx")
prs.save(output_path)
print(f"Sunum olusturuldu: {output_path}")
print(f"Toplam slayt: {len(prs.slides)}")
